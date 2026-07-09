#!/usr/bin/env python3
"""Generate an editable PPTX smoke artifact from the hairy HTML fixture.

This is a prototype artifact generator, not the final browser exporter. It reads
the fixture HTML, extracts slide titles and repeated content patterns, and writes
a native PowerPoint deck with editable text, shapes, chart approximations, and a
table. The purpose is to give the PPTX export plan a concrete verification deck.
"""

from __future__ import annotations

import re
from pathlib import Path

from lxml import html
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
INPUT_HTML = ROOT / "fixtures" / "pptx-export" / "cosmic-canvas-hairy-deck.html"
OUTPUT_PPTX = ROOT / "output" / "pptx" / "cosmic-canvas-hairy-deck-prototype.pptx"

SLIDE_W = 16
SLIDE_H = 9

COLORS = {
    "ink": "172033",
    "muted": "657086",
    "paper": "FBFBF6",
    "blue": "2F6FED",
    "teal": "0E9F93",
    "coral": "F06B58",
    "gold": "E4A72F",
    "plum": "7047A8",
    "lime": "9CC93A",
    "pale_blue": "EAF1FF",
    "pale_teal": "E3F7F4",
    "pale_coral": "FFE5DF",
    "pale_gold": "FFF0C2",
    "white": "FFFFFF",
    "dark": "172033",
}

ACCENTS = [COLORS["blue"], COLORS["teal"], COLORS["coral"], COLORS["gold"], COLORS["plum"], COLORS["lime"]]


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def clean_text(value: str | None) -> str:
    if not value:
        return ""
    return re.sub(r"\s+", " ", value).strip()


def text_content(node) -> str:
    return clean_text(" ".join(node.itertext()))


def first_text(node, xpath: str, fallback: str = "") -> str:
    found = node.xpath(xpath)
    if not found:
        return fallback
    return text_content(found[0])


def class_nodes(node, class_name: str):
    return node.xpath(f".//*[contains(concat(' ', normalize-space(@class), ' '), ' {class_name} ')]")


def set_fill(shape, color: str) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str = "D9E1EC", width: float = 1.0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


def hide_line(shape) -> None:
    shape.line.fill.background()


def add_shape(slide, shape_type, x, y, w, h, fill: str, line: str | None = "D9E1EC"):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    if line is None:
        hide_line(shape)
    else:
        set_line(shape, line)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color: str = COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multiline(slide, lines: list[str], x, y, w, h, size=16, color=COLORS["muted"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.level = 0
        paragraph.space_after = Pt(5)
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return box


def add_card(slide, x, y, w, h, fill=COLORS["white"], line="D9E1EC"):
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line)
    return shape


def add_metric(slide, x, y, w, h, value: str, label: str, accent: str):
    add_card(slide, x, y, w, h, COLORS["white"])
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.12, h, accent, None)
    add_text(slide, value, x + 0.28, y + 0.24, w - 0.42, 0.5, 28, COLORS["ink"], True)
    add_text(slide, label, x + 0.28, y + 0.88, w - 0.42, h - 1.0, 13, COLORS["muted"])


def add_header(slide, source, idx: int, dark: bool = False):
    title = source.get("data-title") or first_text(source, ".//h1|.//h2", f"Slide {idx}")
    eyebrow = first_text(source, ".//*[contains(concat(' ', normalize-space(@class), ' '), ' eyebrow ')]", "Cosmic Canvas")
    color = COLORS["white"] if dark else COLORS["ink"]
    muted = "DDE7F7" if dark else COLORS["muted"]
    add_text(slide, eyebrow.upper(), 0.82, 0.5, 5.3, 0.28, 9, muted, True)
    add_text(slide, title, 0.82, 0.88, 8.8, 0.7, 28, color, True)
    add_text(slide, f"{idx:02d}", 14.65, 0.55, 0.6, 0.3, 12, muted, True, PP_ALIGN.RIGHT)


def add_background(slide, dark: bool = False):
    if dark:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, COLORS["dark"], None)
        add_shape(slide, MSO_SHAPE.OVAL, 11.1, -1.1, 5.2, 5.2, "23568F", None)
        add_shape(slide, MSO_SHAPE.OVAL, -1.3, 4.6, 4.0, 4.0, "6B3F67", None)
    else:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, COLORS["paper"], None)
        add_shape(slide, MSO_SHAPE.OVAL, 13.2, -1.2, 3.5, 3.5, "DCEBFF", None)
        add_shape(slide, MSO_SHAPE.OVAL, -0.8, 6.8, 3.2, 3.2, "E5F7F3", None)


def render_title(slide, source, idx):
    add_background(slide, True)
    add_shape(slide, MSO_SHAPE.PARALLELOGRAM, 10.2, 0.7, 5.9, 6.5, COLORS["blue"], None)
    add_text(slide, "HELIOGRID LABS / OPERATING REVIEW", 0.88, 0.72, 6.0, 0.32, 10, "DDE7F7", True)
    add_text(slide, first_text(source, ".//h1"), 0.88, 1.65, 8.25, 2.45, 47, COLORS["white"], True)
    add_text(slide, first_text(source, ".//p"), 0.92, 4.45, 7.9, 1.2, 17, "DDE7F7")
    chips = [text_content(node) for node in class_nodes(source, "chip")][:3]
    for i, chip in enumerate(chips):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.92 + i * 1.72, 6.15, 1.48, 0.44, "EAF1FF", "B7CFFF")
        add_text(slide, chip, 1.02 + i * 1.72, 6.27, 1.26, 0.18, 8.7, COLORS["blue"], True, PP_ALIGN.CENTER)
    add_card(slide, 11.65, 6.5, 2.6, 1.05, "263B5B", "526C92")
    add_text(slide, "Q3", 11.9, 6.72, 0.86, 0.4, 25, COLORS["white"], True)
    add_text(slide, "Board-ready prototype deck", 12.7, 6.77, 1.24, 0.42, 10, "DDE7F7")


def render_scorecard(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.88, 1.9, 6.8, 5.6)
    add_text(slide, first_text(source, ".//h2"), 1.18, 2.24, 5.9, 1.34, 30, COLORS["ink"], True)
    add_text(slide, first_text(source, ".//p"), 1.18, 3.82, 5.8, 1.0, 15, COLORS["muted"])
    points = [(1.15, 6.45), (2.35, 6.2), (3.35, 6.35), (4.45, 5.72), (5.45, 5.15), (6.3, 5.28), (7.05, 4.78)]
    for i in range(len(points) - 1):
        slide.shapes.add_connector(1, Inches(points[i][0]), Inches(points[i][1]), Inches(points[i + 1][0]), Inches(points[i + 1][1])).line.color.rgb = rgb(COLORS["blue"])
    add_shape(slide, MSO_SHAPE.OVAL, 6.93, 4.66, 0.24, 0.24, COLORS["teal"], None)
    metrics = class_nodes(source, "metric-card")[:4]
    positions = [(8.08, 1.9), (11.58, 1.9), (8.08, 4.82), (11.58, 4.82)]
    for i, (node, (x, y)) in enumerate(zip(metrics, positions)):
        add_metric(slide, x, y, 3.12, 2.38, first_text(node, ".//strong"), first_text(node, ".//span"), ACCENTS[i])


def render_network(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.88, 1.82, 14.25, 6.42, "EDF5FF")
    paths = [((1.55, 3.55), (4.3, 2.2), COLORS["blue"]), ((4.3, 2.2), (7.15, 4.25), COLORS["blue"]), ((7.15, 4.25), (11.7, 3.2), COLORS["teal"]), ((2.55, 6.1), (7.15, 4.25), COLORS["teal"]), ((7.15, 4.25), (12.95, 6.1), COLORS["coral"])]
    for start, end, color in paths:
        connector = slide.shapes.add_connector(1, Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1]))
        connector.line.color.rgb = rgb(color)
        connector.line.width = Pt(4)
    nodes = class_nodes(source, "node")[:6]
    coords = [(1.25, 2.82), (4.0, 2.0), (6.52, 3.65), (11.05, 2.95), (12.32, 5.65), (2.7, 5.9)]
    for i, (node, (x, y)) in enumerate(zip(nodes, coords)):
        add_card(slide, x, y, 1.55, 0.82)
        label = clean_text(node.text)
        small = first_text(node, ".//small")
        add_text(slide, label, x + 0.15, y + 0.14, 0.8, 0.22, 13, COLORS["ink"], True)
        add_text(slide, small, x + 0.15, y + 0.42, 1.1, 0.18, 8.5, COLORS["muted"], True)
        add_shape(slide, MSO_SHAPE.OVAL, x + 1.12, y + 0.22, 0.28, 0.28, ACCENTS[i % len(ACCENTS)], None)


def render_forecast(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.88, 2.0, 7.4, 5.92)
    heights = [0.42, 0.50, 0.57, 0.72, 0.66, 0.81, 0.75, 0.88]
    labels = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"]
    for i, (height, label) in enumerate(zip(heights, labels)):
        x = 1.25 + i * 0.78
        h = height * 3.25
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 6.55 - h, 0.44, h, ACCENTS[i % 3], None)
        add_text(slide, label, x - 0.05, 6.74, 0.54, 0.18, 8.5, COLORS["muted"], True, PP_ALIGN.CENTER)
    add_card(slide, 8.7, 2.0, 6.42, 5.92)
    add_text(slide, "Conversion stressors", 9.05, 2.4, 4.2, 0.4, 22, COLORS["ink"], True)
    add_text(slide, first_text(source, ".//div[contains(@class, 'card')][2]//p"), 9.05, 3.02, 5.4, 1.2, 15, COLORS["muted"])
    chips = [text_content(node) for node in class_nodes(source, "chip")][:4]
    for i, chip in enumerate(chips):
        row, col = divmod(i, 2)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.05 + col * 2.38, 4.85 + row * 0.72, 2.02, 0.42, "EAF1FF", "B7CFFF")
        add_text(slide, chip, 9.16 + col * 2.38, 4.96 + row * 0.72, 1.78, 0.17, 9, COLORS["blue"], True, PP_ALIGN.CENTER)


def render_segments(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.92, 2.1, 4.25, 5.7)
    add_shape(slide, MSO_SHAPE.OVAL, 1.75, 2.74, 2.55, 2.55, COLORS["blue"], None)
    add_shape(slide, MSO_SHAPE.OVAL, 2.18, 3.17, 1.68, 1.68, COLORS["paper"], None)
    add_text(slide, "72%", 2.47, 3.72, 1.1, 0.4, 28, COLORS["ink"], True, PP_ALIGN.CENTER)
    add_text(slide, "Accounts ready for guided automation after simulation replay.", 1.38, 5.72, 3.26, 0.88, 14, COLORS["muted"], False, PP_ALIGN.CENTER)
    metrics = class_nodes(source, "metric-card")[:4]
    positions = [(5.55, 2.1), (10.34, 2.1), (5.55, 5.05), (10.34, 5.05)]
    for i, (node, (x, y)) in enumerate(zip(metrics, positions)):
        add_metric(slide, x, y, 4.25, 2.42, first_text(node, ".//strong"), first_text(node, ".//span"), ACCENTS[i])


def render_timeline(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    x0, y0, gap = 0.98, 3.25, 2.85
    connector = slide.shapes.add_connector(1, Inches(1.35), Inches(3.58), Inches(14.25), Inches(3.58))
    connector.line.color.rgb = rgb(COLORS["teal"])
    connector.line.width = Pt(4)
    for i, item in enumerate(class_nodes(source, "timeline-item")[:5]):
        x = x0 + i * gap
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.2, y0 - 0.08, 0.46, 0.46, COLORS["white"], COLORS["blue"])
        add_card(slide, x, y0 + 0.82, 2.3, 2.38)
        add_text(slide, first_text(item, ".//h3"), x + 0.22, y0 + 1.12, 1.0, 0.28, 17, COLORS["ink"], True)
        add_text(slide, first_text(item, ".//p"), x + 0.22, y0 + 1.62, 1.78, 0.72, 12.5, COLORS["muted"])


def render_architecture(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    headers = [text_content(node) for node in source.xpath(".//div[contains(@class, 'stack')]/h3")]
    modules = [text_content(node) for node in class_nodes(source, "module")]
    for col, header in enumerate(headers[:3]):
        x = 1.0 + col * 4.8
        add_text(slide, header, x, 2.02, 2.5, 0.38, 20, COLORS["ink"], True)
        for row in range(3):
            label = modules[col * 3 + row]
            y = 2.7 + row * 1.24
            add_card(slide, x, y, 3.75, 0.82, ["EAF1FF", "E3F7F4", "FFF0C2"][col])
            add_text(slide, label, x + 0.25, y + 0.25, 2.8, 0.2, 13.5, COLORS["ink"], True)
        if col < 2:
            connector = slide.shapes.add_connector(1, Inches(x + 3.95), Inches(4.2), Inches(x + 4.58), Inches(4.2))
            connector.line.color.rgb = rgb(COLORS["muted"])
            connector.line.width = Pt(2)


def render_waterfall(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.95, 2.02, 14.0, 5.95)
    labels = ["Baseline\n$42M", "Storage\n+8.2", "Dispatch\n+5.4", "Weather\n-3.1", "Automation\n+7.0", "Latency\n-1.5", "Plan\n$58M"]
    heights = [2.0, 1.15, 0.9, 0.65, 1.05, 0.52, 3.35]
    colors = [COLORS["blue"], COLORS["teal"], COLORS["teal"], COLORS["coral"], COLORS["teal"], COLORS["coral"], COLORS["blue"]]
    for i, (label, height, color) in enumerate(zip(labels, heights, colors)):
        x = 1.52 + i * 1.78
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 6.15 - height, 1.05, height, color, None)
        add_text(slide, label, x - 0.18, 6.42, 1.4, 0.48, 9.5, COLORS["muted"], True, PP_ALIGN.CENTER)


def render_heatmap(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    cells = class_nodes(source, "heatcell")
    start_x, start_y = 0.92, 2.15
    col_w, row_h = 2.34, 0.92
    for i, cell in enumerate(cells):
        row, col = divmod(i, 6)
        x, y = start_x + col * col_w, start_y + row * row_h
        klass = cell.get("class", "")
        if "labelcell" in klass:
            fill, text_color, line = COLORS["paper"], COLORS["muted"], None
        elif "high" in klass:
            fill, text_color, line = "FFD8D1", "A02F21", "F4B4AA"
        elif "medium" in klass:
            fill, text_color, line = "FFF0C2", "8A5B00", "E9CC72"
        else:
            fill, text_color, line = "DFF6ED", "0D6F65", "A5E4D2"
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w - 0.1, row_h - 0.12, fill, line)
        add_text(slide, text_content(cell), x + 0.08, y + 0.27, col_w - 0.28, 0.2, 10.5, text_color, True, PP_ALIGN.CENTER)


def render_kanban(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    for col, column in enumerate(class_nodes(source, "column")[:4]):
        x = 0.92 + col * 3.64
        add_card(slide, x, 2.02, 3.22, 5.85, "F5F7FB")
        add_text(slide, first_text(column, ".//h3"), x + 0.25, 2.34, 1.8, 0.3, 18, COLORS["ink"], True)
        for row, task in enumerate(column.xpath(".//*[contains(concat(' ', normalize-space(@class), ' '), ' task ')]")):
            y = 2.9 + row * 1.12
            add_card(slide, x + 0.22, y, 2.68, 0.82, COLORS["white"])
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.22, y, 0.08, 0.82, ACCENTS[col], None)
            add_text(slide, text_content(task), x + 0.42, y + 0.22, 2.1, 0.26, 11.2, COLORS["ink"], True)


def render_matrix(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_card(slide, 0.95, 2.0, 14.0, 5.95, "FFFAF0")
    v = slide.shapes.add_connector(1, Inches(7.95), Inches(2.15), Inches(7.95), Inches(7.72))
    h = slide.shapes.add_connector(1, Inches(1.12), Inches(4.95), Inches(14.72), Inches(4.95))
    for line in (v.line, h.line):
        line.color.rgb = rgb("9AA5B5")
        line.width = Pt(1.5)
    bubbles = [("PDF\nprint", 3.0, 5.7, COLORS["plum"]), ("Hybrid\nPPTX", 6.1, 3.45, COLORS["blue"]), ("Native\ndeck", 10.3, 5.0, COLORS["teal"]), ("Ideal\nfuture", 12.2, 3.05, COLORS["coral"])]
    for text, x, y, color in bubbles:
        add_shape(slide, MSO_SHAPE.OVAL, x, y, 1.18, 1.18, color, None)
        add_text(slide, text, x + 0.18, y + 0.34, 0.82, 0.4, 10.5, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(slide, "lower visual parity", 1.35, 7.55, 2.2, 0.2, 10, COLORS["muted"])
    add_text(slide, "higher visual parity", 12.25, 7.55, 2.2, 0.2, 10, COLORS["muted"], False, PP_ALIGN.RIGHT)
    note = add_text(slide, "higher editability", 14.35, 4.0, 1.35, 0.2, 9, COLORS["muted"], True)
    note.rotation = 90


def render_stress(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    cards = source.xpath(".//*[contains(concat(' ', normalize-space(@class), ' '), ' card ')][.//h3]")[:4]
    positions = [(0.95, 2.35, -2, "EAF1FF"), (5.25, 2.35, 0, "FFE5DF"), (9.55, 2.35, 0, "E3F7F4")]
    for i, (node, (x, y, rotation, fill)) in enumerate(zip(cards[:3], positions)):
        card = add_card(slide, x, y, 3.7, 2.0, fill)
        card.rotation = rotation
        add_text(slide, first_text(node, ".//h3"), x + 0.28, y + 0.28, 2.6, 0.3, 17, COLORS["ink"], True)
        add_text(slide, first_text(node, ".//p"), x + 0.28, y + 0.82, 2.85, 0.6, 12.5, COLORS["muted"])
        if i == 1:
            add_text(slide, "css-clip-path warning", x + 0.28, y + 1.55, 2.2, 0.2, 9.5, COLORS["coral"], True)
        if i == 2:
            add_text(slide, "css-filter warning", x + 0.28, y + 1.55, 2.2, 0.2, 9.5, COLORS["teal"], True)
    add_card(slide, 1.2, 6.05, 13.5, 1.25)
    add_text(slide, "Inline rich text", 1.5, 6.32, 2.6, 0.25, 17, COLORS["ink"], True)
    add_text(slide, "This paragraph includes bold emphasis, italic nuance, and a safe link so the converter can test editable text runs.", 4.0, 6.24, 9.7, 0.38, 13.5, COLORS["muted"])


def render_table(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    rows = source.xpath(".//table[contains(@class, 'export-table')]//tr")
    row_count = len(rows)
    col_count = len(rows[0].xpath("./th|./td")) if rows else 0
    table_shape = slide.shapes.add_table(row_count, col_count, Inches(0.95), Inches(2.2), Inches(14.05), Inches(5.55))
    table = table_shape.table
    widths = [2.4, 4.2, 4.65, 2.2]
    for i, width in enumerate(widths[:col_count]):
        table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        cells = row.xpath("./th|./td")
        for c, cell in enumerate(cells):
            target = table.cell(r, c)
            target.text = text_content(cell)
            fill = target.fill
            fill.solid()
            fill.fore_color.rgb = rgb("EAF1FF" if r == 0 else "FFFFFF")
            for paragraph in target.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(10 if r == 0 else 11.5)
                    run.font.bold = r == 0 or c == 0
                    run.font.color.rgb = rgb(COLORS["ink"] if r else COLORS["blue"])


def render_closing(slide, source, idx):
    add_background(slide, True)
    add_text(slide, "DECISION", 0.88, 0.72, 2.0, 0.28, 10, "DDE7F7", True)
    add_text(slide, first_text(source, ".//h1"), 0.88, 1.42, 10.2, 2.0, 39, COLORS["white"], True)
    add_text(slide, first_text(source, ".//p"), 0.92, 3.8, 8.9, 0.85, 16.5, "DDE7F7")
    metrics = class_nodes(source, "metric-card")[:3]
    for i, node in enumerate(metrics):
        x = 0.95 + i * 4.72
        add_card(slide, x, 5.55, 4.0, 1.65, "263B5B", "526C92")
        add_text(slide, first_text(node, ".//strong"), x + 0.28, 5.88, 0.52, 0.38, 25, COLORS["white"], True)
        add_text(slide, first_text(node, ".//span"), x + 0.95, 5.88, 2.5, 0.42, 11.5, "DDE7F7")


LAYOUTS = {
    "title": render_title,
    "scorecard": render_scorecard,
    "network": render_network,
    "forecast": render_forecast,
    "segments": render_segments,
    "timeline": render_timeline,
    "architecture": render_architecture,
    "waterfall": render_waterfall,
    "heatmap": render_heatmap,
    "kanban": render_kanban,
    "matrix": render_matrix,
    "stress": render_stress,
    "table": render_table,
    "closing": render_closing,
}


def render_fallback(slide, source, idx):
    add_background(slide)
    add_header(slide, source, idx)
    add_text(slide, text_content(source), 0.95, 2.0, 13.5, 5.8, 14, COLORS["muted"])


def build_deck() -> Presentation:
    root = html.fromstring(INPUT_HTML.read_text(encoding="utf-8"))
    slides = root.xpath("//section[contains(concat(' ', normalize-space(@class), ' '), ' slide ')]")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, source in enumerate(slides, start=1):
        ppt_slide = prs.slides.add_slide(blank)
        layout = source.get("data-layout", "")
        render = LAYOUTS.get(layout, render_fallback)
        render(ppt_slide, source, idx)
    return prs


def main() -> None:
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUTPUT_PPTX)
    shape_count = sum(len(slide.shapes) for slide in prs.slides)
    text_shapes = sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    print(f"wrote {OUTPUT_PPTX}")
    print(f"slides={len(prs.slides)} shapes={shape_count} text_shapes={text_shapes}")


if __name__ == "__main__":
    main()

